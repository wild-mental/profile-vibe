"""Fill the Emeritus Educator Bio template (slide3) with ByungJun Park's content."""
import copy
import re
import shutil
import zipfile
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from PIL import Image, ImageDraw

ROOT = "/Users/kidsnote/Workspace/profile-vibe"
BASE = f"{ROOT}/partner-profiles/EMERITUS"
SCRATCH = ("/private/tmp/claude-502/-Users-kidsnote-Workspace-profile-vibe/"
           "861e6084-4c5d-49d4-bdc0-9774d23df0dc/scratchpad")
TEMPLATE = f"{BASE}/Educator Bio Template - 2026 (1).pptx"
OUT = f"{BASE}/output/ByungJun-Park-Emeritus-Bio-2026.pptx"
HEADSHOT_SRC = f"{ROOT}/assets/profile-byungjun-park.png"
HEADSHOT_SQ = f"{SCRATCH}/headshot-square.png"

# Right rail, top to bottom. (source, top_in, width_in) — heights follow the aspect ratio.
# The rail is 10.926->13.335in; buttons occupy 3.039-3.981in and the Emeritus mark 6.906in.
RAIL_LOGOS = [
    (f"{ROOT}/assets/ai-skills-for-everyone-main-logo.jpg", 0.40, 1.98),
    (f"{ROOT}/assets/aztks-skill-introduction-en.jpeg", 1.46, 1.98),
    (f"{ROOT}/assets/vibeworking-consulting-casebook.jpeg", 3.62, 1.98),
]
RAIL_LEFT = 11.18

# Caption above the casebook, matching the layout's deleted "Videos" caption:
# Calibri 12pt bold #00B050, centred. (top, height) in inches.
CAPTION_TEXT = "2026 AX Consulting Casebook"
CAPTION_BOX = (3.13, 0.42)
CAPTION_FONT = ("Calibri", 12, "00B050")

# Link buttons move to the bottom of the rail, above the Emeritus mark at 6.906in.
BUTTON_TOPS = {"Rectangle: Rounded Corners 12": 5.75,
               "Rectangle: Rounded Corners 13": 6.30}

# ---------------------------------------------------------------- content
NAME = "ByungJun Park"

POSITIONS = [
    "Principal AI & IT Education Consultant — Modulabs · UD IMPACT",
    "Former Platform Backend Engineer — Kakao Kidsnote",
    "Former FinTech & Blockchain Backend Engineer — Coinbit Exchange",
]

EXPERIENCE = [
    "A backend engineer turned AI and technology educator. Spent three years building "
    "high-throughput systems for Korean fintech and consumer platforms — sole engineer "
    "responsible for the advertising backend at Kakao Kidsnote, where the platform's ad "
    "revenue more than doubled over two years with zero production incidents after relaunch, "
    "and architect of a real-time alerting system at the Coinbit cryptocurrency exchange "
    "capable of dispatching three million notifications per minute, alongside backend work "
    "for ISMS-P information-security certification. AWS Certified Solutions Architect – Associate.",

    "Delivers in-service technology programs for Hana Financial Group, KT, Kia Motors and "
    "LG HelloVision, and teaches on national workforce tracks at Sahmyook University, "
    "Modulabs and the Korea Information Education Institute — the Sahmyook KDT track "
    "graduated two consecutive cohorts at 100%. Has delivered programs in the United States "
    "and the Middle East, and led a KOICA-sponsored program in Korea for Sri Lankan platform "
    "operators. Contributor to policy research on AI education and curriculum innovation "
    "(KRIVET; Hanbat National University). Teaches in Korean and English.",
]

LEAD_IN = "Industry specific experience as a practitioner and/or educator across globe:"

TABLE = [
    ("Banking & Finance", "Hana Financial Group"),
    ("Digital Assets & Exchanges", "Axiasoft (Coinbit)"),
    ("Internet & Consumer Platform", "Kakao Kidsnote"),
    ("Telecommunications & Media", "KT, LG HelloVision"),
    ("Automotive & Manufacturing", "Kia Motors"),
    ("Government & Public Sector",
     "KOICA, Korea Productivity Center, Korea Software Industry Association"),
    ("Higher Education", "Sahmyook University, Kwangwoon University"),
    ("Corporate & Executive Training",
     "Samsung Multicampus, Modulabs, Goorm EDU, Sparta Coding Club"),
]

CORE_TOPICS = ("AX Productivity | AI Skills & Harness Engineering"
               " | Cloud Modernization in Finance")

SUB_TOPICS = ("Spec-Driven Development for Both Engineers and Non-engineers | Hands-off AI "
              "Workflow Automation with Maximum Job Ownership | Prompt and Context Engineering "
              "| Microservices and Kubernetes Operations | Kafka Event Streaming "
              "| Ad-Tech Targeting and ML Monetization")

EDUCATION = [
    "MS in AI & Big Data · MBA — aSSIST University & SDG MS (Geneva, Switzerland)",
    "BA in Computer Science — Korea National Open University",
    "BA in Political Science & Communication — Sogang University",
]

LOCATION = "Seoul, South Korea"

# 3 template buttons; only two verified public links exist -> third is deleted.
VIDEO_LINKS = [
    ("Rectangle: Rounded Corners 12", "Portfolio", "https://pbjworking.com"),
    ("Rectangle: Rounded Corners 13", "LinkedIn", "https://www.linkedin.com/in/pbjworking/"),
]
VIDEO_DELETE = ["Rectangle: Rounded Corners 14"]

TABLE_TOP = Inches(4.68)          # template ships 4.868 -> table overflows the 7.5in slide

# ---------------------------------------------------------------- marker fallback
# slideLayout9's location-pin ("Graphic 14") ships as an SVG-only <a:blip>: no r:embed
# raster, just the asvg:svgBlip extension. Windows PowerPoint reads it; macOS/Keynote/
# Google Slides/LibreOffice draw an empty frame. We rasterise ppt/media/image11.svg and
# attach it as the raster fallback, exactly as PowerPoint itself stores SVG pictures.
MARKER_MEDIA = "ppt/media/image20.png"
MARKER_RID = "rId4"
MARKER_FILL = (0, 190, 109, 255)          # .MsftOfcThm_Accent1_Fill_v2 -> #00BE6D
MARKER_VIEWBOX = 96

# image11.svg, transcribed subpath by subpath. Opposite winding -> the disc is a hole.
MARKER_DISC = [
    ("M", (48, 44)),
    ("C", (43, 44, 39, 40, 39, 35)),
    ("C", (39, 30, 43, 26, 48, 26)),
    ("C", (53, 26, 57, 30, 57, 35)),
    ("C", (57, 40, 53, 44, 48, 44)),
]
MARKER_TEARDROP = [
    ("M", (48, 14)),
    ("C", (41.1, 14, 34.6, 17.4, 30.7, 23.2)),
    ("C", (26.8, 28.9, 26, 36.2, 28.5, 42.7)),
    ("L", (38, 63.7)),
    ("L", (46.2, 80.9)),
    ("C", (46.5, 81.6, 47.2, 82, 48, 82)),
    ("C", (48.8, 82, 49.5, 81.6, 49.8, 80.9)),
    ("L", (58, 63.7)),
    ("L", (67.5, 42.7)),
    ("C", (70, 36.2, 69.2, 28.9, 65.3, 23.2)),
    ("C", (61.4, 17.4, 54.9, 14, 48, 14)),
]


def flatten(subpath, scale, steps=64):
    """Turn an SVG subpath (M / L / cubic C) into a polygon in device space."""
    pts, cur = [], None
    for cmd, args in subpath:
        if cmd == "M":
            cur = args
            pts.append(cur)
        elif cmd == "L":
            cur = args
            pts.append(cur)
        elif cmd == "C":
            x0, y0 = cur
            x1, y1, x2, y2, x3, y3 = args
            for i in range(1, steps + 1):
                t = i / steps
                u = 1 - t
                pts.append((
                    u**3 * x0 + 3 * u*u*t * x1 + 3 * u*t*t * x2 + t**3 * x3,
                    u**3 * y0 + 3 * u*u*t * y1 + 3 * u*t*t * y2 + t**3 * y3,
                ))
            cur = (x3, y3)
    return [(x * scale, y * scale) for x, y in pts]


def render_marker(path, px=512, ss=4):
    """Rasterise the marker to a transparent PNG (supersampled, then downscaled)."""
    size = px * ss
    scale = size / MARKER_VIEWBOX
    img = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    ImageDraw.Draw(img).polygon(flatten(MARKER_TEARDROP, scale), fill=MARKER_FILL)
    hole = Image.new("L", (size, size), 255)
    ImageDraw.Draw(hole).polygon(flatten(MARKER_DISC, scale), fill=0)
    img.putalpha(Image.composite(img.getchannel("A"),
                                 Image.new("L", (size, size), 0), hole))
    img.resize((px, px), Image.LANCZOS).save(path)

# ---------------------------------------------------------------- helpers
A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def set_paragraphs(shape, lines):
    """Replace a text frame's content with `lines`, preserving paragraph 0's formatting."""
    tf = shape.text_frame
    paras = tf.paragraphs
    template_p = copy.deepcopy(paras[0]._p)

    for p in list(paras)[1:]:
        p._p.getparent().remove(p._p)

    def write(para, text):
        runs = para.runs
        if not runs:
            para.add_run().text = text
            return
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)

    write(paras[0], lines[0])
    anchor = paras[0]._p
    for _ in lines[1:]:
        new_p = copy.deepcopy(template_p)
        anchor.addnext(new_p)
        anchor = new_p
    for para, line in zip(tf.paragraphs[1:], lines[1:]):
        write(para, line)


def fill_cell(cell, text):
    """Write text into a table cell, inheriting the cell's own end-paragraph run props."""
    para = cell.text_frame.paragraphs[0]
    for r in list(para.runs):
        r._r.getparent().remove(r._r)
    run = para.add_run()
    run.text = text
    end_rpr = para._p.find(f"{A}endParaRPr")
    if end_rpr is not None:
        rpr = copy.deepcopy(end_rpr)
        rpr.tag = f"{A}rPr"
        old = run._r.find(f"{A}rPr")
        if old is not None:
            run._r.remove(old)
        run._r.insert(0, rpr)


# ---------------------------------------------------------------- images
Image.open(HEADSHOT_SRC).convert("RGB").crop((300, 80, 1620, 1400)) \
     .resize((1000, 1000), Image.LANCZOS).save(HEADSHOT_SQ)


def prepare_logo(src, index, max_px=1200):
    """Downscale a rail logo so the deck stays small; returns (path, aspect ratio)."""
    im = Image.open(src).convert("RGB")
    if im.size[0] > max_px:
        im = im.resize((max_px, round(max_px * im.size[1] / im.size[0])), Image.LANCZOS)
    out = f"{SCRATCH}/rail-logo-{index}.jpg"
    im.save(out, quality=88)
    return out, im.size[0] / im.size[1]

# ---------------------------------------------------------------- build
prs = Presentation(TEMPLATE)

sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst)[:2]:            # keep only slide 3
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)
assert len(prs.slides) == 1

slide = prs.slides[0]
by_name = {sh.name: sh for sh in slide.shapes}

set_paragraphs(by_name["Text Placeholder 15"], [NAME])
set_paragraphs(by_name["Text Placeholder 10"], POSITIONS)
set_paragraphs(by_name["Text Placeholder 9"], EXPERIENCE + [LEAD_IN])
set_paragraphs(by_name["Text Placeholder 6"], [CORE_TOPICS])
set_paragraphs(by_name["Text Placeholder 7"], [SUB_TOPICS])
set_paragraphs(by_name["Text Placeholder 4"], EDUCATION)
set_paragraphs(by_name["Text Placeholder 5"], [LOCATION])

tbl_shape = by_name["Table 2"]
tbl_shape.top = TABLE_TOP
table = tbl_shape.table
for i, (industry, companies) in enumerate(TABLE, start=1):
    fill_cell(table.cell(i, 0), industry)
    fill_cell(table.cell(i, 1), companies)

for name, label, url in VIDEO_LINKS:
    btn = by_name[name]
    set_paragraphs(btn, [label])
    btn.click_action.hyperlink.address = url
    btn.top = Inches(BUTTON_TOPS[name])
for name in VIDEO_DELETE:
    el = by_name[name]._element
    el.getparent().remove(el)

by_name["Picture Placeholder 3"].insert_picture(HEADSHOT_SQ)

for i, (src, top, width_in) in enumerate(RAIL_LOGOS):
    path, ratio = prepare_logo(src, i)
    w = Inches(width_in)
    slide.shapes.add_picture(path, Inches(RAIL_LEFT), Inches(top),
                             width=w, height=int(w / ratio))

cap_top, cap_h = CAPTION_BOX
face, pts, colour = CAPTION_FONT
caption = slide.shapes.add_textbox(Inches(RAIL_LEFT), Inches(cap_top),
                                   Inches(RAIL_LOGOS[0][2]), Inches(cap_h))
caption.name = "Casebook Caption"
tf = caption.text_frame
tf.word_wrap = True
tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para = tf.paragraphs[0]
para.alignment = PP_ALIGN.CENTER
run = para.add_run()
run.text = CAPTION_TEXT
run.font.name = face
run.font.size = Pt(pts)
run.font.bold = True
run.font.color.rgb = RGBColor.from_string(colour)

prs.save(OUT)

# ------------------------------------------- metadata + marker raster fallback
MARKER_PNG = f"{SCRATCH}/marker.png"
render_marker(MARKER_PNG)

tmp = f"{SCRATCH}/_meta.pptx"
with zipfile.ZipFile(OUT) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.infolist():
        data = zin.read(item.filename)
        if item.filename == "ppt/slideLayouts/slideLayout9.xml":
            x = data.decode("utf-8")
            i = x.find('name="Graphic 14"')
            assert i != -1, "Graphic 14 not found in slideLayout9"
            j = x.find("<a:blip>", i)
            assert j != -1, "Graphic 14 has no <a:blip>"
            x = x[:j] + f'<a:blip r:embed="{MARKER_RID}">' + x[j + len("<a:blip>"):]

            # drop the "Videos" caption above the link buttons
            k = x.find('name="TextBox 1"')
            assert k != -1, '"TextBox 1" not found in slideLayout9'
            start = x.rfind("<p:sp>", 0, k)
            end = x.find("</p:sp>", k) + len("</p:sp>")
            assert "Videos" in x[start:end], '"TextBox 1" is not the Videos caption'
            x = x[:start] + x[end:]
            data = x.encode("utf-8")
        elif item.filename == "ppt/slideLayouts/_rels/slideLayout9.xml.rels":
            x = data.decode("utf-8")
            assert f'Id="{MARKER_RID}"' not in x
            rel = (f'<Relationship Id="{MARKER_RID}" Type="http://schemas.openxmlformats.org'
                   f'/officeDocument/2006/relationships/image" '
                   f'Target="../media/{MARKER_MEDIA.split("/")[-1]}"/>')
            x = x.replace("</Relationships>", rel + "</Relationships>")
            data = x.encode("utf-8")
        elif item.filename == "docProps/app.xml":
            x = data.decode("utf-8")
            x = x.replace("<Slides>3</Slides>", "<Slides>1</Slides>")
            x = re.sub(r"<vt:lpstr>Guidelines for bio template[^<]*</vt:lpstr>", "", x)
            x = x.replace("<vt:lpstr>PowerPoint Presentation</vt:lpstr>"
                          "<vt:lpstr>PowerPoint Presentation</vt:lpstr>",
                          "<vt:lpstr>ByungJun Park - Educator Bio</vt:lpstr>")
            x = x.replace('<vt:variant><vt:i4>3</vt:i4></vt:variant>'
                          '</vt:vector></HeadingPairs>',
                          '<vt:variant><vt:i4>1</vt:i4></vt:variant>'
                          '</vt:vector></HeadingPairs>')
            x = x.replace('<vt:vector size="9" baseType="lpstr">',
                          '<vt:vector size="7" baseType="lpstr">')
            data = x.encode("utf-8")
        elif item.filename == "docProps/core.xml":
            x = data.decode("utf-8")
            x = re.sub(r"<dc:title>[^<]*</dc:title>",
                       "<dc:title>ByungJun Park - Educator Bio 2026</dc:title>", x)
            data = x.encode("utf-8")
        zout.writestr(item, data)
    zout.write(MARKER_PNG, MARKER_MEDIA)
shutil.move(tmp, OUT)
print("saved:", OUT)
